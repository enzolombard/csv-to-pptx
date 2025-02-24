import csv
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from copy import deepcopy
import re

# Default values
PPT_TEMPLATE_PATH = "pptx_template.pptx"  # Replace with your default template path
IMAGE_TO_INSERT = "tacel-logo.png"  # Replace with your default image path

def main(csv_file_path, output_ppt_path):
    # Load the PowerPoint template
    try:
        presentation = Presentation(PPT_TEMPLATE_PATH)
    except FileNotFoundError:
        print(f"Error: PowerPoint template file '{PPT_TEMPLATE_PATH}' not found.")
        sys.exit(1)

    # Read data from the CSV file
    try:
        with open(csv_file_path, 'r') as csvfile:
            reader = csv.DictReader(csvfile)
            rows = list(reader)  # Store all rows for multiple slide generation
    except FileNotFoundError:
        print(f"Error: CSV file '{csv_file_path}' not found.")
        sys.exit(1)

    if not rows:
        print("Error: CSV file is empty or not properly formatted.")
        sys.exit(1)

    # Access the first slide as the template
    if len(presentation.slides) == 0:
        print("Error: PowerPoint template does not contain any slides.")
        sys.exit(1)

    template_slide = presentation.slides[0]

    # Function to duplicate a slide
    def duplicate_slide(presentation, slide):
        """Duplicate a slide, including all shapes and content."""
        slide_index = presentation.slides.index(slide)
        slide_layout = slide.slide_layout
        new_slide = presentation.slides.add_slide(slide_layout)
        
        for shape in slide.shapes:
            new_element = deepcopy(shape._element)
            new_slide.shapes._spTree.insert(len(new_slide.shapes._spTree), new_element)
        return new_slide
    
    # Function to clean strings
    def clean_text(text):
        return re.sub(r"^\d+\.\s*", "", text)  # Removes leading number, dot, and optional space
    
    # Function to format SECTION TITLE
    def format_section_title(section_text, steps_text):
        """
        Extracts and formats the SECTION TITLE based on:
        "SECTION X - STEP #
        -----------------------------------
        STEP TITLE"
        """

        global LAST_SECTION_NUMBER 

        # Extract SECTION number
        section_match = re.match(r"(\d+)\.\s*(.+)", section_text)  # Extracts "1" and "Service Box Setup"
        if section_match:
            section_number = section_match.group(1)  # Extract number
            LAST_SECTION_NUMBER = section_number  # Update last known section number
        else:
            section_number = LAST_SECTION_NUMBER  # Use last known valid section num

        # Extract STEP number (removes "S" from "S1")
        step_match = re.match(r"S(\d+)", steps_text)  # Finds "S1" and extracts "1"
        step_number = step_match.group(1) if step_match else "X"  # Default "X" if not found

        # Extract STEP TITLE from text in parentheses
        title_match = re.search(r"\((.*?)\)", steps_text)  # Extracts text inside parentheses
        step_title = title_match.group(1).upper() if title_match else steps_text.split()[0]  # Default first word if no match

        # Format final string
        formatted_text = f"SECTION {section_number} - STEP {step_number}\n-----------------------------------\n{step_title}"
        return formatted_text
    
    def format_head(section_text, steps_text):
        """
        Extracts and formats the HEAD section as:
        "(SECTION TITLE) SEC X STEP #"
        """

        global LAST_SECTION_TITLE, LAST_SECTION_NUMBER  # Use last known section values

        # Extract SECTION number and title
        section_match = re.match(r"(\d+)\.\s*(.+)", section_text)  # Extracts "2" and "Wire Prep"
        if section_match:
            section_number = section_match.group(1)  # Extract section number
            section_title = section_match.group(2).upper()  # Extract and convert section title to uppercase
            LAST_SECTION_NUMBER = section_number  # Update last known section number
            LAST_SECTION_TITLE = section_title  # Update last known section title
        else:
            section_number = LAST_SECTION_NUMBER  # Use last known valid section number
            section_title = LAST_SECTION_TITLE  # Use last known valid section title

        # Extract STEP number (removes "S" from "S1")
        step_match = re.match(r"S(\d+)", steps_text)
        step_number = step_match.group(1) if step_match else "X"  # Default "X" if not found

        # Format final string
        formatted_text = f"({section_title}) SEC {section_number} STEP {step_number}"

        return formatted_text
    

    def format_steps(steps_text):
        """
        Extracts the instructional part of STEPS by removing:
        - "S1 (Step Title)" → Removes "S1" and anything inside parentheses.
        - Leaves only the main instruction.
        
        Example:
        Input:  "S1 (Cutting 6 AWG Wires)  Cut the 6 AWG Green wire to 4.5 inches..."
        Output: "Cut the 6 AWG Green wire to 4.5 inches..."
        """

        # Remove "S1 " or any step number format "S#"
        steps_text = re.sub(r"^S\d+\s*", "", steps_text)

        # Remove the step title in parentheses
        steps_text = re.sub(r"\(.*?\)\s*", "", steps_text)

        return steps_text.strip()  # Trim leading/trailing spaces

    # Function to delete default title and text boxes
    def remove_default_shapes(slide):
        """Remove default title and text boxes from the slide."""
        shapes_to_remove = []
        for shape in slide.shapes:
            # Identify title and text box placeholders
            if shape.is_placeholder and shape.placeholder_format.type in ['TITLE', 'BODY']:
                shapes_to_remove.append(shape)

        for shape in shapes_to_remove:
            spTree = shape._element.getparent()
            spTree.remove(shape._element)
            print(f"Removed shape '{shape.name}' from slide.")

    # Function to replace pictures on a slide with a specific image
    def replace_pictures(slide, image_path):
        """Replace all pictures in a slide with a specified image."""
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:  # Check if the shape is a picture
                # Get position and size of the current picture
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                # Delete the existing picture
                spTree = shape._element.getparent()
                spTree.remove(shape._element)

                # Insert the new image in the same position
                slide.shapes.add_picture(image_path, left, top, width, height)
                print(f"Replaced picture at position ({left}, {top}) with image '{image_path}'")

    # Create slides for each row in the CSV
    for row_idx, row in enumerate(rows):
        print(f"Processing row {row_idx + 1}: {row}")

        # Duplicate the template slide
        new_slide = duplicate_slide(presentation, template_slide)

        # Remove default title and text boxes
        remove_default_shapes(new_slide)


        # Replace placeholders with CSV data
        for shape in new_slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = run.text  # Get current placeholder text
                        
                        if "{SECTION TITLE}" in text:
                            processed_value = format_section_title(row["SECTION"], row["STEPS"])
                            print(f"Replacing {text} with {processed_value}")  # Debugging
                            run.text = run.text.replace("{SECTION TITLE}", processed_value)

                        elif "{HEAD}" in text:
                            processed_value = format_head(row["SECTION"], row["STEPS"])  # Custom function to process {Head}
                            print(f"Replacing {text} with {processed_value}")  # Debugging
                            run.text = run.text.replace("{HEAD}", processed_value)

                        elif "{STEPS}" in text:
                            processed_value = format_steps(row["STEPS"])  # Custom function to process {Steps}
                            print(f"Replacing {text} with {processed_value}")  # Debugging
                            run.text = run.text.replace("{STEPS}", processed_value)

        # Replace all pictures in the slide
        replace_pictures(new_slide, IMAGE_TO_INSERT)

    # Save the updated PowerPoint presentation
    if len(presentation.slides) > 0:
        presentation.save(output_ppt_path)
        print(f"PowerPoint presentation saved to {output_ppt_path}")
    else:
        print("No slides were generated.")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python script.py <csv_file_path> <output_ppt_path>")
        sys.exit(1)

    csv_file_path = sys.argv[1]
    output_ppt_path = sys.argv[2]

    main(csv_file_path, output_ppt_path)